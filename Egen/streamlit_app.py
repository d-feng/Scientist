import streamlit as st
import asyncio
import json
import pandas as pd
from datetime import datetime
import plotly.express as px
import plotly.graph_objects as go
from typing import Dict, List, Any
import os
import sys
import PyPDF2
import pptx

# Import the orchestrator (assumes the file is in the same directory)
try:
    from ai_agent_orchestrator_2 import AIAgentOrchestrator, LLMProvider, AgentType, AnalysisDepth
    ORCHESTRATOR_AVAILABLE = True
except ImportError:
    ORCHESTRATOR_AVAILABLE = False
    st.error("ai_agent_orchestrator_2.py not found. Please ensure it's in the same directory as this app.")

# Configure page
st.set_page_config(
    page_title="Red Team AI Agent",
    page_icon="🔬",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for better styling
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        color: #1f77b4;
        text-align: center;
        margin-bottom: 2rem;
    }
    .agent-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 1rem;
        border-radius: 10px;
        margin: 0.5rem 0;
        color: white;
    }
    .metric-card {
        background: #f8f9fa;
        padding: 1rem;
        border-radius: 8px;
        border-left: 4px solid #1f77b4;
    }
    .issue-high {
        background: #ffebee;
        border-left: 4px solid #f44336;
        padding: 0.5rem;
        margin: 0.25rem 0;
    }
    .issue-medium {
        background: #fff3e0;
        border-left: 4px solid #ff9800;
        padding: 0.5rem;
        margin: 0.25rem 0;
    }
    .issue-low {
        background: #f3e5f5;
        border-left: 4px solid #9c27b0;
        padding: 0.5rem;
        margin: 0.25rem 0;
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state
if 'orchestrator' not in st.session_state:
    st.session_state.orchestrator = None
if 'analysis_results' not in st.session_state:
    st.session_state.analysis_results = None
if 'analysis_history' not in st.session_state:
    st.session_state.analysis_history = []

# Helper functions
def initialize_orchestrator(provider: str) -> bool:
    """Initialize the AI orchestrator"""
    try:
        llm_provider = LLMProvider.OPENAI if provider == "OpenAI" else LLMProvider.ANTHROPIC
        st.session_state.orchestrator = AIAgentOrchestrator(llm_provider=llm_provider)
        return True
    except Exception as e:
        st.error(f"Failed to initialize orchestrator: {str(e)}")
        return False

async def run_analysis(slide_content: str, slide_type: str, agent_types: List[str], analysis_depth: str) -> Dict[str, Any]:
    """Run the analysis using the orchestrator"""
    if not st.session_state.orchestrator:
        raise ValueError("Orchestrator not initialized")
    
    await st.session_state.orchestrator.initialize()
    
    results = await st.session_state.orchestrator.analyze_slide(
        slide_content=slide_content,
        slide_type=slide_type,
        agent_types=agent_types,
        analysis_depth=analysis_depth
    )
    
    return results

def display_analysis_results(results: Dict[str, Any]):
    """Display analysis results in a structured format"""
    if not results:
        st.warning("No analysis results to display.")
        return
    
    # Overall metrics
    st.subheader("📊 Analysis Overview")
    
    col1, col2, col3, col4 = st.columns(4)
    
    total_issues = 0
    total_recommendations = 0
    avg_confidence = 0
    agent_count = 0
    
    for agent_name, result in results.items():
        if agent_name != "cross_validation" and "error" not in result:
            total_issues += len(result.get("technical_issues", []))
            total_recommendations += len(result.get("recommendations", []))
            avg_confidence += result.get("confidence_score", 0)
            agent_count += 1
    
    if agent_count > 0:
        avg_confidence /= agent_count
    
    with col1:
        st.metric("Total Issues", total_issues)
    with col2:
        st.metric("Total Recommendations", total_recommendations)
    with col3:
        st.metric("Average Confidence", f"{avg_confidence:.2f}")
    with col4:
        st.metric("Agents Used", agent_count)
    
    # Agent-specific results
    st.subheader("🤖 Agent Analysis Results")
    
    for agent_name, result in results.items():
        if agent_name == "cross_validation":
            continue
            
        if "error" in result:
            st.error(f"**{agent_name.title()} Agent**: {result['error']}")
            continue
        
        with st.expander(f"🔍 {agent_name.title()} Agent Results", expanded=True):
            
            # Agent metrics
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Issues Found", len(result.get("technical_issues", [])))
            with col2:
                st.metric("Recommendations", len(result.get("recommendations", [])))
            with col3:
                st.metric("Confidence", f"{result.get('confidence_score', 0):.2f}")
            
            # Vague language issues
            vague_issues = result.get("vague_language", [])
            if vague_issues:
                st.write("**🔤 Vague Language Issues:**")
                for issue in vague_issues:
                    severity_class = f"issue-{issue.get('severity', 'low')}"
                    st.markdown(f"""
                    <div class="{severity_class}">
                        <strong>Text:</strong> "{issue.get('text', '')}"<br>
                        <strong>Category:</strong> {issue.get('category', '')}<br>
                        <strong>Suggestion:</strong> {issue.get('suggestion', '')}<br>
                        <strong>Severity:</strong> {issue.get('severity', '')} (Confidence: {issue.get('confidence', 0):.2f})
                    </div>
                    """, unsafe_allow_html=True)
            
            # Technical issues
            tech_issues = result.get("technical_issues", [])
            if tech_issues:
                st.write("**⚠️ Technical Issues:**")
                for issue in tech_issues:
                    severity_class = f"issue-{issue.get('severity', 'low')}"
                    st.markdown(f"""
                    <div class="{severity_class}">
                        <strong>Issue:</strong> {issue.get('issue', '')}<br>
                        <strong>Category:</strong> {issue.get('category', '')}<br>
                        <strong>Recommendation:</strong> {issue.get('recommendation', '')}<br>
                        <strong>Severity:</strong> {issue.get('severity', '')} (Confidence: {issue.get('confidence', 0):.2f})
                    </div>
                    """, unsafe_allow_html=True)
            
            # Recommendations
            recommendations = result.get("recommendations", [])
            if recommendations:
                st.write("**💡 Recommendations:**")
                for rec in recommendations:
                    priority_color = {
                        "critical": "#f44336",
                        "high": "#ff9800",
                        "medium": "#2196f3",
                        "low": "#4caf50"
                    }.get(rec.get("priority", "low"), "#4caf50")
                    
                    st.markdown(f"""
                    <div style="border-left: 4px solid {priority_color}; padding: 0.5rem; margin: 0.25rem 0; background: #f8f9fa;">
                        <strong>Type:</strong> {rec.get('type', '')}<br>
                        <strong>Recommendation:</strong> {rec.get('recommendation', '')}<br>
                        <strong>Priority:</strong> {rec.get('priority', '')} (Confidence: {rec.get('confidence', 0):.2f})<br>
                        <strong>Rationale:</strong> {rec.get('rationale', '')}
                    </div>
                    """, unsafe_allow_html=True)
    
    # Cross-validation results
    if "cross_validation" in results:
        st.subheader("🔄 Cross-Validation Analysis")
        cross_val = results["cross_validation"]
        
        if "error" not in cross_val:
            col1, col2 = st.columns(2)
            
            with col1:
                st.metric("Consistency Score", f"{cross_val.get('consistency_score', 0):.2f}")
            
            with col2:
                consensus_count = len(cross_val.get("consensus_issues", []))
                st.metric("Consensus Issues", consensus_count)
            
            # Consensus issues
            consensus_issues = cross_val.get("consensus_issues", [])
            if consensus_issues:
                st.write("**✅ Consensus Issues:**")
                for issue in consensus_issues:
                    st.markdown(f"""
                    <div class="issue-{issue.get('severity', 'low')}">
                        <strong>Category:</strong> {issue.get('category', '')}<br>
                        <strong>Description:</strong> {issue.get('description', '')}<br>
                        <strong>Agent Agreement:</strong> {issue.get('agent_count', 0)} agents<br>
                        <strong>Severity:</strong> {issue.get('severity', '')}
                    </div>
                    """, unsafe_allow_html=True)
            
            # Conflicting findings
            conflicts = cross_val.get("conflicting_findings", [])
            if conflicts:
                st.write("**⚔️ Conflicting Findings:**")
                for conflict in conflicts:
                    st.markdown(f"""
                    <div style="border-left: 4px solid #ff5722; padding: 0.5rem; margin: 0.25rem 0; background: #fff3e0;">
                        <strong>Type:</strong> {conflict.get('type', '')}<br>
                        <strong>Conflict:</strong> {conflict.get('conflict', '')}<br>
                        <strong>Agents Involved:</strong> {', '.join(conflict.get('agents_involved', []))}<br>
                        <strong>Resolution:</strong> {conflict.get('resolution', '')}
                    </div>
                    """, unsafe_allow_html=True)
            
            # Priority recommendations
            priority_recs = cross_val.get("priority_recommendations", [])
            if priority_recs:
                st.write("**🎯 Priority Recommendations:**")
                for rec in priority_recs:
                    priority_color = {
                        "critical": "#f44336",
                        "high": "#ff9800",
                        "medium": "#2196f3",
                        "low": "#4caf50"
                    }.get(rec.get("priority", "low"), "#4caf50")
                    
                    st.markdown(f"""
                    <div style="border-left: 4px solid {priority_color}; padding: 0.5rem; margin: 0.25rem 0; background: #f8f9fa;">
                        <strong>Priority:</strong> {rec.get('priority', '')}<br>
                        <strong>Recommendation:</strong> {rec.get('recommendation', '')}<br>
                        <strong>Rationale:</strong> {rec.get('rationale', '')}<br>
                        <strong>Supporting Agents:</strong> {', '.join(rec.get('supporting_agents', []))}
                    </div>
                    """, unsafe_allow_html=True)
        else:
            st.error(f"Cross-validation failed: {cross_val.get('error', 'Unknown error')}")

def create_analysis_dashboard(results: Dict[str, Any]):
    """Create visualization dashboard for analysis results"""
    if not results:
        return
    
    st.subheader("📈 Analysis Dashboard")
    
    # Prepare data for visualization
    agent_data = []
    issue_data = []
    recommendation_data = []
    
    for agent_name, result in results.items():
        if agent_name == "cross_validation" or "error" in result:
            continue
            
        # Agent performance data
        agent_data.append({
            "Agent": agent_name.title(),
            "Issues": len(result.get("technical_issues", [])),
            "Recommendations": len(result.get("recommendations", [])),
            "Confidence": result.get("confidence_score", 0),
            "Processing Time": result.get("processing_time", 0)
        })
        
        # Issue severity distribution
        for issue in result.get("technical_issues", []):
            issue_data.append({
                "Agent": agent_name.title(),
                "Severity": issue.get("severity", "low"),
                "Category": issue.get("category", "other")
            })
        
        # Recommendation priority distribution
        for rec in result.get("recommendations", []):
            recommendation_data.append({
                "Agent": agent_name.title(),
                "Priority": rec.get("priority", "low"),
                "Type": rec.get("type", "other")
            })
    
    if agent_data:
        df_agents = pd.DataFrame(agent_data)
        
        col1, col2 = st.columns(2)
        
        with col1:
            # Agent performance comparison
            fig = px.bar(df_agents, x="Agent", y=["Issues", "Recommendations"], 
                        title="Agent Performance Comparison",
                        barmode="group")
            st.plotly_chart(fig, use_container_width=True)
        
        with col2:
            # Confidence scores
            fig = px.scatter(df_agents, x="Agent", y="Confidence", 
                           size="Issues", color="Processing Time",
                           title="Agent Confidence vs Performance")
            st.plotly_chart(fig, use_container_width=True)
        
        if issue_data:
            df_issues = pd.DataFrame(issue_data)
            
            col1, col2 = st.columns(2)
            
            with col1:
                # Issue severity distribution
                fig = px.histogram(df_issues, x="Severity", color="Agent",
                                 title="Issue Severity Distribution")
                st.plotly_chart(fig, use_container_width=True)
            
            with col2:
                # Issue category breakdown
                fig = px.pie(df_issues, names="Category", title="Issue Categories")
                st.plotly_chart(fig, use_container_width=True)

# Main app
def main():
    st.markdown('<h1 class="main-header">🔬 Red Team AI Agent</h1>', unsafe_allow_html=True)
    st.markdown("### Analyze scientific slide content with specialized AI agents")
    
    if not ORCHESTRATOR_AVAILABLE:
        st.stop()
    
    # Sidebar configuration
    st.sidebar.header("⚙️ Configuration")
    
    # LLM Provider selection
    llm_provider = st.sidebar.selectbox(
        "Choose LLM Provider:",
        ["OpenAI", "Anthropic"],
        help="Select the language model provider"
    )
    
    # Initialize orchestrator
    if not st.session_state.orchestrator:
        if initialize_orchestrator(llm_provider):
            st.sidebar.success("✅ Orchestrator initialized!")
        else:
            st.sidebar.error("❌ Failed to initialize orchestrator")
    
    # Analysis configuration
    st.sidebar.header("🔍 Analysis Settings")
    
    slide_type = st.sidebar.selectbox(
        "Slide Type:",
        ["general", "compound_profile", "mechanism", "clinical_data", "safety_profile", "competitive_analysis"],
        help="Select the type of slide content"
    )
    
    agent_types = st.sidebar.multiselect(
        "Select Agents:",
        ["structural", "pharmacokinetic", "clinical", "all"],
        default=["all"],
        help="Choose which specialized agents to use"
    )
    
    analysis_depth = st.sidebar.selectbox(
        "Analysis Depth:",
        ["standard", "detailed", "comprehensive"],
        index=1,
        help="Select the depth of analysis"
    )
    
    # Main content area
    st.header("📝 Slide Content Input")
    
    # Content input methods
    input_method = st.radio(
        "Choose input method:",
        ["Text Input", "File Upload", "PDF Upload", "PowerPoint Upload", "Example Content"],
        horizontal=True
    )
    
    slide_content = ""
    
    if input_method == "Text Input":
        slide_content = st.text_area(
            "Enter slide content:",
            height=200,
            placeholder="Paste your slide content here..."
        )
    
    elif input_method == "File Upload":
        uploaded_file = st.file_uploader(
            "Upload slide content file:",
            type=['txt', 'md']
        )
        if uploaded_file:
            slide_content = uploaded_file.read().decode('utf-8')
            st.text_area("File content:", value=slide_content, height=200, disabled=True)

    elif input_method == "PDF Upload":
        uploaded_pdf = st.file_uploader(
            "Upload presentation or poster PDF:",
            type=['pdf']
        )
        if uploaded_pdf:
            try:
                pdf_reader = PyPDF2.PdfReader(uploaded_pdf)
                text_content = ""
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text_content += page.extract_text() + "\n\n"
                slide_content = text_content
                st.text_area("Extracted PDF content:", value=slide_content, height=200, disabled=True)
            except Exception as e:
                st.error(f"Error reading PDF: {e}")

    elif input_method == "PowerPoint Upload":
        uploaded_pptx = st.file_uploader(
            "Upload presentation PowerPoint file:",
            type=['pptx']
        )
        if uploaded_pptx:
            try:
                from pptx import Presentation
                prs = Presentation(uploaded_pptx)
                text_content = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content += shape.text + "\n"
                slide_content = text_content
                st.text_area("Extracted PowerPoint content:", value=slide_content, height=200, disabled=True)
            except Exception as e:
                st.error(f"Error reading PowerPoint: {e}")

    elif input_method == "Example Content":
        example_content = {
            "Compound Profile": """
            Our novel compound XYZ-123 shows promising binding affinity to the target protein.
            The compound exhibits good bioavailability and minimal side effects in preclinical studies.
            Initial pharmacokinetic studies suggest once-daily dosing potential.
            The mechanism involves selective inhibition of the target enzyme.
            Safety profile appears favorable with no major toxicities observed.
            """,
            "Clinical Data": """
            Phase II clinical trial results demonstrate significant efficacy in patient populations.
            Primary endpoint was met with statistical significance (p<0.05).
            Secondary endpoints showed consistent improvements across all measured parameters.
            Safety profile was acceptable with manageable adverse events.
            Patient reported outcomes were positive with high satisfaction scores.
            """,
            "Mechanism of Action": """
            The compound acts through selective binding to the active site of the target enzyme.
            Crystal structure analysis reveals key interactions with critical amino acid residues.
            Inhibition kinetics demonstrate competitive inhibition with high selectivity.
            Downstream effects include modulation of key signaling pathways.
            The mechanism is well-characterized and supports the therapeutic approach.
            """
        }
        
        selected_example = st.selectbox("Select example:", list(example_content.keys()))
        slide_content = example_content[selected_example]
        st.text_area("Example content:", value=slide_content, height=200, disabled=True)
    
    # Analysis button
    if st.button("🚀 Run Analysis", type="primary", disabled=not (slide_content and st.session_state.orchestrator)):
        if not slide_content:
            st.warning("Please provide slide content to analyze.")
        elif not st.session_state.orchestrator:
            st.warning("Please initialize the orchestrator first.")
        else:
            with st.spinner("Analyzing slide content..."):
                try:
                    # Run analysis
                    results = asyncio.run(run_analysis(
                        slide_content=slide_content,
                        slide_type=slide_type,
                        agent_types=agent_types,
                        analysis_depth=analysis_depth
                    ))
                    
                    # Store results
                    st.session_state.analysis_results = results
                    
                    # Add to history
                    st.session_state.analysis_history.append({
                        "timestamp": datetime.now(),
                        "slide_type": slide_type,
                        "agent_types": agent_types,
                        "analysis_depth": analysis_depth,
                        "results": results
                    })
                    
                    st.success("✅ Analysis completed successfully!")
                    
                except Exception as e:
                    st.error(f"❌ Analysis failed: {str(e)}")
    
    # Display results
    if st.session_state.analysis_results:
        st.header("📊 Analysis Results")
        
        # Create tabs for different views
        tab1, tab2, tab3 = st.tabs(["📋 Detailed Results", "📈 Dashboard", "📑 Export"])
        
        with tab1:
            display_analysis_results(st.session_state.analysis_results)
        
        with tab2:
            create_analysis_dashboard(st.session_state.analysis_results)
        
        with tab3:
            st.subheader("📁 Export Results")
            
            # JSON export
            json_data = json.dumps(st.session_state.analysis_results, indent=2, default=str)
            st.download_button(
                label="📥 Download JSON",
                data=json_data,
                file_name=f"analysis_results_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json",
                mime="application/json"
            )
            
            # CSV export (flattened data)
            try:
                # Create flattened data for CSV
                csv_data = []
                for agent_name, result in st.session_state.analysis_results.items():
                    if agent_name != "cross_validation" and "error" not in result:
                        for issue in result.get("technical_issues", []):
                            csv_data.append({
                                "Agent": agent_name,
                                "Type": "Technical Issue",
                                "Content": issue.get("issue", ""),
                                "Category": issue.get("category", ""),
                                "Severity": issue.get("severity", ""),
                                "Confidence": issue.get("confidence", 0)
                            })
                        for rec in result.get("recommendations", []):
                            csv_data.append({
                                "Agent": agent_name,
                                "Type": "Recommendation",
                                "Content": rec.get("recommendation", ""),
                                "Category": rec.get("type", ""),
                                "Severity": rec.get("priority", ""),
                                "Confidence": rec.get("confidence", 0)
                            })
                
                if csv_data:
                    df = pd.DataFrame(csv_data)
                    csv_string = df.to_csv(index=False)
                    st.download_button(
                        label="📥 Download CSV",
                        data=csv_string,
                        file_name=f"analysis_results_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
                        mime="text/csv"
                    )
            except Exception as e:
                st.error(f"CSV export failed: {str(e)}")
    
    # Analysis history
    if st.session_state.analysis_history:
        st.header("📚 Analysis History")
        
        with st.expander("View previous analyses"):
            for i, analysis in enumerate(reversed(st.session_state.analysis_history)):
                st.write(f"**Analysis {len(st.session_state.analysis_history) - i}** - {analysis['timestamp'].strftime('%Y-%m-%d %H:%M:%S')}")
                st.write(f"Slide Type: {analysis['slide_type']}, Agents: {', '.join(analysis['agent_types'])}, Depth: {analysis['analysis_depth']}")
                
                if st.button(f"Load Analysis {len(st.session_state.analysis_history) - i}", key=f"load_{i}"):
                    st.session_state.analysis_results = analysis['results']
                    st.rerun()
                
                st.divider()

if __name__ == "__main__":
    main()